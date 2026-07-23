"""
Generate a readable Dashboard Guide PPT — 2 slides, standard widescreen,
large fonts. All 9 tabs covered with descriptions.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── colours ───────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x11, 0x2B)
BG_CARD  = RGBColor(0x1A, 0x1D, 0x4C)
PINK     = RGBColor(0xEC, 0x48, 0x99)
YELLOW   = RGBColor(0xFA, 0xCC, 0x15)
ORANGE   = RGBColor(0xFB, 0x92, 0x3C)
CYAN     = RGBColor(0x22, 0xD3, 0xEE)
GREEN    = RGBColor(0x34, 0xD3, 0x99)
RED      = RGBColor(0xF8, 0x71, 0x71)
VIOLET   = RGBColor(0xA7, 0x8B, 0xFA)
BLUE     = RGBColor(0x60, 0xA5, 0xFA)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x9C, 0xA3, 0xAF)
DGRAY    = RGBColor(0x4B, 0x55, 0x63)
LGRAY    = RGBColor(0x6B, 0x72, 0x80)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── helpers ───────────────────────────────────────────────────────────
def _bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def _box(sl, l, t, w, h, txt, sz=14, clr=WHITE, bold=False, al=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz)
    p.font.color.rgb = clr; p.font.bold = bold; p.font.name = 'Calibri'; p.alignment = al

def _ml(sl, l, t, w, h, lines, sz=11):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        txt, c = item[0], item[1] if len(item) > 1 else GRAY
        b = item[2] if len(item) > 2 else False
        s = item[3] if len(item) > 3 else sz
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(s); p.font.color.rgb = c
        p.font.bold = b; p.font.name = 'Calibri'; p.space_after = Pt(2)

def _card(sl, l, t, w, h, border=PINK):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = BG_CARD
    s.line.color.rgb = border; s.line.width = Pt(1.5)
    s.shadow.inherit = False

def _rect(sl, l, t, w, h, clr):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()

# =====================================================================
# SLIDE 1 — Tabs 1-5  (Overview, Behavioral, Barriers, Needs, Segments)
# =====================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s1)

# Header
_rect(s1, 0, 0, 13.333, 0.85, BG_CARD)
_box(s1, 0.4, 0.08, 9, 0.4, "Blinkit Discovery Engine -- Dashboard Guide", 22, YELLOW, True)
_box(s1, 0.4, 0.48, 12, 0.3,
     "Next.js 14 | React 18 | Recharts | FastAPI | MongoDB Atlas | PostgreSQL | Vercel | Render | GitHub Actions",
     10, LGRAY)
_box(s1, 11, 0.12, 2, 0.3, "Slide 1 / 2", 10, DGRAY, False, PP_ALIGN.RIGHT)

# ── Tab 1: Overview ──────────────────────────────────────────────────
_card(s1, 0.3, 1.0, 4.05, 3.05, CYAN)
_ml(s1, 0.45, 1.05, 3.8, 2.9, [
    ("Tab 1: Overview", CYAN, True, 14),
    ("", WHITE),
    ("4 KPI Cards: Total Reviews (with 7-day trend %),", WHITE, False, 11),
    ("Avg Rating, Issues Detected (NLP), Barriers.", WHITE, False, 11),
    ("", WHITE),
    ("Charts:", CYAN, True, 11),
    ("- 30-day Area Chart (reviews + avg rating)", GRAY, False, 10),
    ("- Rating Distribution (1-5 star bars)", GRAY, False, 10),
    ("- Platform Distribution (Play / iOS / Web)", GRAY, False, 10),
    ("- Scatter: Exploration vs Satisfaction", GRAY, False, 10),
    ("- Grouped Bar: Current / New / Churned users", GRAY, False, 10),
    ("- Heatmap: Issue intensity by category x time", GRAY, False, 10),
    ("- 3 Insight Cards: Top Barrier, Need, Frustration", GRAY, False, 10),
])

# ── Tab 2: Behavioral Analysis ──────────────────────────────────────
_card(s1, 4.55, 1.0, 4.05, 3.05, BLUE)
_ml(s1, 4.7, 1.05, 3.8, 2.9, [
    ("Tab 2: Behavioral Analysis", BLUE, True, 14),
    ("", WHITE),
    ("30-day Line Chart: daily review volume.", WHITE, False, 11),
    ("Scatter: Exploration (X) vs Satisfaction (Y),", WHITE, False, 11),
    ("bubble size = Order Frequency, color = segment.", WHITE, False, 11),
    ("", WHITE),
    ("Sections:", BLUE, True, 11),
    ("- Grouped Bar: users per category", GRAY, False, 10),
    ("- Segment x Frustration Crosstab (matrix)", GRAY, False, 10),
    ("- Recent Activity Feed (scrollable reviews)", GRAY, False, 10),
    ("- Q&A Cards: Habit Impact, Discovery Methods", GRAY, False, 10),
    ("- TrendAnalysis: segment + volume trends", GRAY, False, 10),
])

# ── Tab 3: Barriers ──────────────────────────────────────────────────
_card(s1, 8.8, 1.0, 4.25, 3.05, RED)
_ml(s1, 8.95, 1.05, 4.0, 2.9, [
    ("Tab 3: Barriers", RED, True, 14),
    ("", WHITE),
    ("4 KPIs: Avg Rating (color-coded), Barrier Count,", WHITE, False, 11),
    ("Frustration Mentions (% of reviews), Unmet Needs.", WHITE, False, 11),
    ("", WHITE),
    ("Charts:", RED, True, 11),
    ("- Frustration Frequency (horizontal bars, NLP)", GRAY, False, 10),
    ("- Barrier Severity Chart (Recharts, deduplicated)", GRAY, False, 10),
    ("- Rating Distribution (5-star breakdown)", GRAY, False, 10),
    ("- User Segment Impact (exposure level)", GRAY, False, 10),
    ("- AI-generated Key Barriers Insight", GRAY, False, 10),
    ("- TrendAnalysis: top 5 frustration trends", GRAY, False, 10),
])

# ── Tab 4: Unmet Needs ───────────────────────────────────────────────
_card(s1, 0.3, 4.2, 6.35, 3.1, ORANGE)
_ml(s1, 0.45, 4.25, 6.1, 3.0, [
    ("Tab 4: Unmet Needs", ORANGE, True, 14),
    ("", WHITE),
    ("Opportunity Matrix (8-row table):", WHITE, False, 11),
    ("Columns: Area | Frequency | Impact/10 | Effort | RICE Score | Status", WHITE, False, 10),
    ("Top: Faster Delivery ETAs (RICE 87, Critical), Better Search (79), Real-time Tracking (72).", GRAY, False, 10),
    ("", WHITE),
    ("4 KPI Cards: Reviews Analyzed, Avg Rating, Unmet Needs Found, Issues Reported.", GRAY, False, 10),
    ("Sentiment-Impact Quadrant: 2x2 bubble chart (Quick Wins / Critical / Monitor / Vocal Minority).", GRAY, False, 10),
    ("Top Feature Requests: Schedule deliveries (89), Loyalty points (76), Stock availability (64).", GRAY, False, 10),
    ("Competitive Gap: vs Zepto (delivery speed), BigBasket (range), Swiggy Instamart (UX).", GRAY, False, 10),
    ("TrendAnalysis: Reviews, Rating, top 4 frustration frequency trends.", GRAY, False, 10),
])

# ── Tab 5: User Segments ─────────────────────────────────────────────
_card(s1, 6.85, 4.2, 6.2, 3.1, GREEN)
_ml(s1, 7.0, 4.25, 5.9, 3.0, [
    ("Tab 5: User Segments", GREEN, True, 14),
    ("", WHITE),
    ("Scatter Plot: Exploration Score (X) vs Purchase Frequency (Y),", WHITE, False, 11),
    ("bubble = Order Value, 3 color-coded segment groups.", WHITE, False, 11),
    ("", WHITE),
    ("Grouped Bar: High / Medium / Low exploration across categories.", GRAY, False, 10),
    ("Q&A Cards: Experimental Segments (AI), Repeat Purchase Behavior (AI).", GRAY, False, 10),
    ("", WHITE),
    ("Trend Analysis:", GREEN, True, 10),
    ("  High Exploration Users: +11.1%  |  Medium: -4.8%  |  Low: 0% stable", GRAY, False, 10),
    ("Data: behavioral_analysis.user_segments, question_answers.*", LGRAY, False, 9),
])

# =====================================================================
# SLIDE 2 — Tabs 6-9  (Discovery, AI Insights, Feedback, Architecture)
# =====================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s2)

# Header
_rect(s2, 0, 0, 13.333, 0.85, BG_CARD)
_box(s2, 0.4, 0.08, 9, 0.4, "Blinkit Discovery Engine -- Dashboard Guide (cont.)", 22, YELLOW, True)
_box(s2, 0.4, 0.48, 12, 0.3,
     "Data Flow: Play Store -> GitHub Actions (6h) -> FastAPI NLP -> MongoDB + PostgreSQL -> Next.js -> Recharts",
     10, CYAN)
_box(s2, 11, 0.12, 2, 0.3, "Slide 2 / 2", 10, DGRAY, False, PP_ALIGN.RIGHT)

# ── Tab 6: Discovery Patterns ────────────────────────────────────────
_card(s2, 0.3, 1.0, 4.05, 3.05, VIOLET)
_ml(s2, 0.45, 1.05, 3.8, 2.9, [
    ("Tab 6: Discovery Patterns", VIOLET, True, 14),
    ("", WHITE),
    ("30-day Line Chart: discovery activity over time.", WHITE, False, 11),
    ("", WHITE),
    ("Sections:", VIOLET, True, 11),
    ("- Heatmap: categories (delivery, quality, app,", GRAY, False, 10),
    ("  support, pricing) x time slots", GRAY, False, 10),
    ("- Q&A: Information Needs Before Exploration", GRAY, False, 10),
    ("- Q&A: Discovery Methods (AI-generated)", GRAY, False, 10),
    ("", WHITE),
    ("Trend Analysis:", VIOLET, True, 10),
    ("  High Exploration % (+15.7%), Reviews This Week,", GRAY, False, 10),
    ("  Categories Covered, Avg Rating.", GRAY, False, 10),
])

# ── Tab 7: AI Insights ───────────────────────────────────────────────
_card(s2, 4.55, 1.0, 4.05, 3.05, PINK)
_ml(s2, 4.7, 1.05, 3.8, 2.9, [
    ("Tab 7: AI Insights", PINK, True, 14),
    ("", WHITE),
    ("30-day Area Chart: insight generation volume.", WHITE, False, 11),
    ("Scatter: Sentiment (X) vs Impact (Y), bubble = Freq.", WHITE, False, 11),
    ("", WHITE),
    ("Sections:", PINK, True, 11),
    ("- Grouped Bar: Delivery, Quality, App, Support,", GRAY, False, 10),
    ("  Pricing, Variety categories", GRAY, False, 10),
    ("- Frustration Cards (2-col grid): theme name,", GRAY, False, 10),
    ("  impact badge, frequency count", GRAY, False, 10),
    ("- Recurring Frustrations (AI text summary)", GRAY, False, 10),
    ("- Trend: Frustrations, Barriers, Unmet Needs", GRAY, False, 10),
])

# ── Tab 8: User Feedback ─────────────────────────────────────────────
_card(s2, 8.8, 1.0, 4.25, 3.05, CYAN)
_ml(s2, 8.95, 1.05, 4.0, 2.9, [
    ("Tab 8: User Feedback", CYAN, True, 14),
    ("", WHITE),
    ("30-day Line Chart: feedback submission volume.", WHITE, False, 11),
    ("", WHITE),
    ("Sections:", CYAN, True, 11),
    ("- Recent Feedback Feed (scrollable review list)", GRAY, False, 10),
    ("  Each card: text, platform, star rating, date", GRAY, False, 10),
    ("- Heatmap: Feedback Intensity by Category & Time", GRAY, False, 10),
    ("", WHITE),
    ("Trend Analysis:", CYAN, True, 10),
    ("  Total Feedback (week trend %)", GRAY, False, 10),
    ("  Positive % (+5.3%), Negative % (-9.1%)", GRAY, False, 10),
    ("  Avg Rating: current vs previous period", GRAY, False, 10),
])

# ── Tab 9: Architecture ──────────────────────────────────────────────
_card(s2, 0.3, 4.2, 8.5, 3.1, YELLOW)
_ml(s2, 0.45, 4.25, 8.2, 3.0, [
    ("Tab 9: Architecture (11 sections embedded in dashboard)", YELLOW, True, 14),
    ("", WHITE),
    ("1. Production System Overview: Data Sources > GitHub Actions Scheduler > FastAPI Backend > MongoDB + PostgreSQL > Next.js Frontend", WHITE, False, 10),
    ("2. Data Flow Pipeline: Ingestion > Processing > API Serving > CI/CD Scheduling", WHITE, False, 10),
    ("3. Real-Time Computation Pipeline: 12-step NLP process cards + segment & frustration classification rules", WHITE, False, 10),
    ("4. Backend API Endpoints: 9 routes with HTTP method badges (GET /api/v1/data/dashboard, POST /ingest/trigger, etc.)", WHITE, False, 10),
    ("5. Database Schema: MongoDB document schema + PostgreSQL tables (data_sources, categories, barriers, unmet_needs) + Redis cache", WHITE, False, 10),
    ("6. Infrastructure & Cost: 6 KPI stats, free-tier cost breakdown, performance benchmarks, system limitations", WHITE, False, 10),
    ("7. Security: MongoDB Atlas auth, HTTPS everywhere, CORS whitelist, JWT tokens, network isolation", WHITE, False, 10),
    ("8. Tech Stack: Frontend (10 libs) + Backend (10 libs) + Infrastructure (8 services) with version numbers", GRAY, False, 10),
    ("9. Repo Structure: Full file tree for backend/ and frontend/ directories", GRAY, False, 10),
    ("10. End-to-End Request Flow: 10 numbered steps from GitHub Actions cron trigger to rendered Recharts", GRAY, False, 10),
    ("11. Environment Config: Render + Vercel environment variables (MONGODB_URI, DATABASE_URL, NEXT_PUBLIC_API_URL)", GRAY, False, 10),
])

# ── Components + Tech Summary ────────────────────────────────────────
_card(s2, 9.0, 4.2, 4.05, 3.1, BLUE)
_ml(s2, 9.15, 4.25, 3.8, 3.0, [
    ("Reusable Components", BLUE, True, 14),
    ("", WHITE),
    ("StatCard -- KPI with icon, value, trend", WHITE, False, 10),
    ("InsightCard -- Summary with metric", WHITE, False, 10),
    ("TimeSeriesChart -- Line / Area (Recharts)", WHITE, False, 10),
    ("RankedBarChart -- Horizontal bar chart", WHITE, False, 10),
    ("ScatterPlot -- Bubble scatter", WHITE, False, 10),
    ("GroupedBarChart -- Multi-metric bars", WHITE, False, 10),
    ("HeatmapChart -- Category x Time", WHITE, False, 10),
    ("TrendAnalysis -- Metric comparison", WHITE, False, 10),
    ("FilterPanel -- Global filters", WHITE, False, 10),
    ("ArchitectureDiagram -- Full docs", WHITE, False, 10),
    ("", WHITE),
    ("Global Filters: Platform | Rating | Sentiment", LGRAY, False, 9),
    ("| Category | Date Range (all proportional)", LGRAY, False, 9),
])

# ── save ──────────────────────────────────────────────────────────────
output_path = r"c:\Top fellow - blinkit\Blinkit_Dashboard_Guide.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
