"""
2-slide PPT: Blinkit Customer Problem Report + PM-backed Solutions
Slide 1 — 3 major problems (from dashboard + web data)
Slide 2 — 3 solutions with PM frameworks
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── colours ───────────────────────────────────────────────────────────
BG      = RGBColor(0x0F, 0x11, 0x2B)
BG_CARD = RGBColor(0x1A, 0x1D, 0x4C)
PINK    = RGBColor(0xEC, 0x48, 0x99)
YELLOW  = RGBColor(0xFA, 0xCC, 0x15)
ORANGE  = RGBColor(0xFB, 0x92, 0x3C)
CYAN    = RGBColor(0x22, 0xD3, 0xEE)
GREEN   = RGBColor(0x34, 0xD3, 0x99)
RED     = RGBColor(0xF8, 0x71, 0x71)
VIOLET  = RGBColor(0xA7, 0x8B, 0xFA)
BLUE    = RGBColor(0x60, 0xA5, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
DGRAY   = RGBColor(0x4B, 0x55, 0x63)
LGRAY   = RGBColor(0x6B, 0x72, 0x80)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── helpers ───────────────────────────────────────────────────────────
def _bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def _box(sl, l, t, w, h, txt, sz=14, clr=WHITE, bold=False, al=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz)
    p.font.color.rgb = clr; p.font.bold = bold; p.font.name = 'Calibri'; p.alignment = al
    return tb

def _ml(sl, l, t, w, h, lines, sz=11):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        txt, c = item[0], item[1] if len(item) > 1 else GRAY
        b = item[2] if len(item) > 2 else False
        s = item[3] if len(item) > 3 else sz
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(s); p.font.color.rgb = c
        p.font.bold = b; p.font.name = 'Calibri'; p.space_after = Pt(1)

def _card(sl, l, t, w, h, border=PINK):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = BG_CARD
    s.line.color.rgb = border; s.line.width = Pt(1.5)
    s.shadow.inherit = False

def _rect(sl, l, t, w, h, clr):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()

def _severity_bar(sl, l, t, w, pct, clr):
    """Tiny progress bar."""
    _rect(sl, l, t, w, 0.12, DGRAY)
    _rect(sl, l, t, w * pct, 0.12, clr)

# =====================================================================
# SLIDE 1 — PROBLEM IDENTIFICATION
# =====================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s1)

# Header banner
_rect(s1, 0, 0, 13.333, 1.0, BG_CARD)
_box(s1, 0.4, 0.1, 9, 0.45, "Blinkit Customer Problem Report", 24, YELLOW, True)
_box(s1, 0.4, 0.55, 12, 0.35,
     "Sources: Dashboard NLP Analysis (barriers, frustrations, segments, RICE scores)  +  Web Sentiment (CurlyTales, ZeeNews, Udit Vani, Reddit, X/Twitter)",
     8.5, LGRAY)
_box(s1, 11.5, 0.15, 1.6, 0.3, "1 / 2", 10, DGRAY, False, PP_ALIGN.RIGHT)

# ── PROBLEM 1 ────────────────────────────────────────────────────────
_card(s1, 0.3, 1.15, 4.1, 3.0, RED)
_ml(s1, 0.5, 1.22, 3.75, 2.85, [
    ("PROBLEM 1", RED, True, 9),
    ("Delayed & Unreliable Deliveries", RED, True, 14),
    ("", WHITE),
    ("Dashboard Evidence:", WHITE, True, 10),
    ("  Barrier severity: 9.2 / 10  (highest)", GRAY, False, 9),
    ("  RICE Score: 87  --  Status: CRITICAL", GRAY, False, 9),
    ("  342 review mentions (top opportunity area)", GRAY, False, 9),
    ("  Delivery Speed: Blinkit 72% vs Zepto 85%", GRAY, False, 9),
    ("  Sentiment Quadrant: HIGH impact + HIGH negative", GRAY, False, 9),
    ("  NLP frustration theme: 'late delivery' = HIGH impact", GRAY, False, 9),
    ("", WHITE),
    ("Web Evidence:", LGRAY, True, 10),
    ("  10-min promise -> 30-45 min in Tier-2 (Udit Vani)", LGRAY, False, 9),
    ("  Riders penalized for traffic/weather delays", LGRAY, False, 9),
    ("  Missing items, wrong items post-rush delivery", LGRAY, False, 9),
])
# severity bar
_box(s1, 0.5, 3.82, 1.2, 0.18, "Severity", 7, DGRAY)
_severity_bar(s1, 1.5, 3.86, 2.7, 0.92, RED)
_box(s1, 3.55, 3.82, 0.6, 0.18, "9.2/10", 7, RED, True)

# ── PROBLEM 2 ────────────────────────────────────────────────────────
_card(s1, 4.6, 1.15, 4.1, 3.0, ORANGE)
_ml(s1, 4.8, 1.22, 3.75, 2.85, [
    ("PROBLEM 2", ORANGE, True, 9),
    ("Poor Refund & Customer Support", ORANGE, True, 14),
    ("", WHITE),
    ("Dashboard Evidence:", WHITE, True, 10),
    ("  Barrier severity: 8.8 / 10  (2nd highest)", GRAY, False, 9),
    ("  RICE Score: 61  --  Status: HIGH", GRAY, False, 9),
    ("  134 mentions of 'Better Refund Process'", GRAY, False, 9),
    ("  Customer Support: Blinkit 58% vs Competitor 70%", GRAY, False, 9),
    ("  Refund bubble: HIGH impact in Sentiment Quadrant", GRAY, False, 9),
    ("  Frustration themes: refund, support = HIGH impact", GRAY, False, 9),
    ("", WHITE),
    ("Web Evidence:", LGRAY, True, 10),
    ("  No refund for rotten eggs (CurlyTales, Nov 2024)", LGRAY, False, 9),
    ("  20-min return window too short (Economic Times)", LGRAY, False, 9),
    ("  Auto-responses only, no live agent (Reddit)", LGRAY, False, 9),
])
_box(s1, 4.8, 3.82, 1.2, 0.18, "Severity", 7, DGRAY)
_severity_bar(s1, 5.8, 3.86, 2.7, 0.88, ORANGE)
_box(s1, 7.85, 3.82, 0.6, 0.18, "8.8/10", 7, ORANGE, True)

# ── PROBLEM 3 ────────────────────────────────────────────────────────
_card(s1, 8.9, 1.15, 4.1, 3.0, VIOLET)
_ml(s1, 9.1, 1.22, 3.75, 2.85, [
    ("PROBLEM 3", VIOLET, True, 9),
    ("Low New Category Exploration", VIOLET, True, 14),
    ("", WHITE),
    ("Dashboard Evidence:", WHITE, True, 10),
    ("  Medium Exploration users declining: -4.8%", GRAY, False, 9),
    ("  Low Exploration segment: STABLE (no growth)", GRAY, False, 9),
    ("  Only High Exploration growing (+11.1%)", GRAY, False, 9),
    ("  'Category Recommendations' RICE: 55 (Low)", GRAY, False, 9),
    ("  Discovery tab: users stick to known categories", GRAY, False, 9),
    ("  Q&A: 'barriers to exploration' = top AI insight", GRAY, False, 9),
    ("", WHITE),
    ("Web / PM Evidence:", LGRAY, True, 10),
    ("  Head vs tail category conflict (Deepak Krishnan)", LGRAY, False, 9),
    ("  ~49% users buy only for immediate needs (NextLeap)", LGRAY, False, 9),
    ("  No personalized discovery path in current UX", LGRAY, False, 9),
])
_box(s1, 9.1, 3.82, 1.2, 0.18, "Severity", 7, DGRAY)
_severity_bar(s1, 10.1, 3.86, 2.7, 0.65, VIOLET)
_box(s1, 12.15, 3.82, 0.6, 0.18, "6.5/10", 7, VIOLET, True)

# ── BOTTOM SECTION: Impact Summary ───────────────────────────────────
_rect(s1, 0.3, 4.35, 12.7, 0.02, DGRAY)

_card(s1, 0.3, 4.55, 6.2, 2.7, CYAN)
_ml(s1, 0.5, 4.6, 5.9, 2.6, [
    ("IMPACT SUMMARY", CYAN, True, 12),
    ("", WHITE),
    ("Business Impact:", WHITE, True, 10),
    ("  - Avg Rating: 3.7 / 5 (below 4.0 threshold for app store visibility)", GRAY, False, 9),
    ("  - Frustration mentions cover ~15-20% of all reviews analyzed", GRAY, False, 9),
    ("  - Delivery + Refund barriers alone account for 476+ negative mentions", GRAY, False, 9),
    ("  - Medium/Low exploration users = majority, but NOT growing", GRAY, False, 9),
    ("  - Competitive gap: losing to Zepto (speed), BigBasket (range), Swiggy (UX)", GRAY, False, 9),
    ("", WHITE),
    ("User Retention Risk:", YELLOW, True, 10),
    ("  Users who hit delivery/refund barriers churn 2-3x faster (industry avg).", GRAY, False, 9),
    ("  Low exploration = low AOV = lower lifetime value per customer.", GRAY, False, 9),
])

_card(s1, 6.7, 4.55, 6.3, 2.7, YELLOW)
_ml(s1, 6.9, 4.6, 5.9, 2.6, [
    ("DASHBOARD KPIs AT A GLANCE", YELLOW, True, 12),
    ("", WHITE),
    ("Avg Rating ............................... 3.7 / 5", WHITE, False, 10),
    ("Total Barriers Tracked ................... 15+", WHITE, False, 10),
    ("Frustration Themes (NLP) ................. 8+", WHITE, False, 10),
    ("Unmet Needs Identified ................... 12+", WHITE, False, 10),
    ("Top Opportunity (RICE) ................... Faster ETAs (87)", WHITE, False, 10),
    ("Delivery Gap vs Zepto .................... -13%", RED, False, 10),
    ("Support Gap vs Competitors ............... -12%", RED, False, 10),
    ("Medium Exploration Trend ................. -4.8% (declining)", ORANGE, False, 10),
    ("High Exploration Trend ................... +11.1% (but small %)", GREEN, False, 10),
])

# =====================================================================
# SLIDE 2 — SOLUTIONS
# =====================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s2)

# Header
_rect(s2, 0, 0, 13.333, 1.0, BG_CARD)
_box(s2, 0.4, 0.1, 9, 0.45, "3 PM-Backed Solutions for Blinkit", 24, YELLOW, True)
_box(s2, 0.4, 0.55, 12, 0.35,
     "Frameworks: RICE Prioritization  |  JTBD (Jobs To Be Done)  |  Zepto Dark Store Model  |  Smart Basket (NextLeap PM Case Study)  |  Swiggy Separate Fulfillment",
     8.5, LGRAY)
_box(s2, 11.5, 0.15, 1.6, 0.3, "2 / 2", 10, DGRAY, False, PP_ALIGN.RIGHT)

# ── SOLUTION 1 ────────────────────────────────────────────────────────
_card(s2, 0.3, 1.15, 4.1, 4.05, GREEN)
_ml(s2, 0.5, 1.22, 3.75, 3.9, [
    ("SOLUTION 1  --  Solves Problem 1", GREEN, True, 9),
    ("Predictive ETA + Smart Routing Engine", GREEN, True, 13),
    ("", WHITE),
    ("What:", WHITE, True, 10),
    ("  AI-powered dynamic ETA that factors in real-time", GRAY, False, 9),
    ("  traffic, weather, dark store load, and rider density.", GRAY, False, 9),
    ("  Show honest ETA ranges (e.g. '12-18 min') instead", GRAY, False, 9),
    ("  of aspirational '10 min'. Under-promise, over-deliver.", GRAY, False, 9),
    ("", WHITE),
    ("PM Framework -- JTBD:", CYAN, True, 10),
    ("  WHEN I order urgently, I want to KNOW exactly when", CYAN, False, 9),
    ("  it will arrive SO I can plan accordingly.", CYAN, False, 9),
    ("", WHITE),
    ("How (from Zepto's playbook):", WHITE, True, 10),
    ("  1. Real-time demand heatmaps per dark store zone", GRAY, False, 9),
    ("  2. Rider allocation AI: match closest + fastest rider", GRAY, False, 9),
    ("  3. Proactive delay notification if ETA will slip >5 min", GRAY, False, 9),
    ("  4. Tier-2 city: increase dark store density gradually", GRAY, False, 9),
    ("", WHITE),
    ("Expected Impact:", YELLOW, True, 10),
    ("  RICE: 87 -> resolves top-1 critical barrier", YELLOW, False, 9),
    ("  Target: reduce delivery complaints by 40-50%", YELLOW, False, 9),
    ("  Competitive: close 13% gap vs Zepto", YELLOW, False, 9),
])

# ── SOLUTION 2 ────────────────────────────────────────────────────────
_card(s2, 4.6, 1.15, 4.1, 4.05, ORANGE)
_ml(s2, 4.8, 1.22, 3.75, 3.9, [
    ("SOLUTION 2  --  Solves Problem 2", ORANGE, True, 9),
    ("Instant Resolution Engine + Extended Returns", ORANGE, True, 13),
    ("", WHITE),
    ("What:", WHITE, True, 10),
    ("  Replace bot-only support with tiered resolution:", GRAY, False, 9),
    ("  Auto-refund for items <Rs.200 (no questions asked).", GRAY, False, 9),
    ("  Extend return window from 20 min to 2 hours for", GRAY, False, 9),
    ("  quality issues. Live agent escalation within 90 sec.", GRAY, False, 9),
    ("", WHITE),
    ("PM Framework -- RICE:", CYAN, True, 10),
    ("  Reach: 134+ mentions. Impact: 8.8/10.", CYAN, False, 9),
    ("  Confidence: HIGH (web data confirms). Effort: MEDIUM.", CYAN, False, 9),
    ("", WHITE),
    ("How:", WHITE, True, 10),
    ("  1. Photo-based auto-approval for damaged items", GRAY, False, 9),
    ("  2. Instant wallet credit (not coupons with expiry)", GRAY, False, 9),
    ("  3. 'Trust Score' for repeat users: auto-approve refunds", GRAY, False, 9),
    ("  4. 24/7 live chat with <90 sec first response SLA", GRAY, False, 9),
    ("", WHITE),
    ("Expected Impact:", YELLOW, True, 10),
    ("  Close 12% customer support gap vs competitors", YELLOW, False, 9),
    ("  Target: 60% reduction in negative refund mentions", YELLOW, False, 9),
    ("  Improve avg rating from 3.7 -> 4.0+ within 2 quarters", YELLOW, False, 9),
])

# ── SOLUTION 3 ────────────────────────────────────────────────────────
_card(s2, 8.9, 1.15, 4.1, 4.05, VIOLET)
_ml(s2, 9.1, 1.22, 3.75, 3.9, [
    ("SOLUTION 3  --  Solves Problem 3", VIOLET, True, 9),
    ("AI Discovery Engine + Smart Basket", VIOLET, True, 13),
    ("", WHITE),
    ("What:", WHITE, True, 10),
    ("  Personalized category discovery to convert Medium/Low", GRAY, False, 9),
    ("  exploration users into High exploration. AI-curated", GRAY, False, 9),
    ("  'Discover New' carousel + Smart Basket auto-restock", GRAY, False, 9),
    ("  with bundled items from unexplored categories.", GRAY, False, 9),
    ("", WHITE),
    ("PM Framework -- JTBD:", CYAN, True, 10),
    ("  WHEN I reorder my usual items, HELP ME discover", CYAN, False, 9),
    ("  relevant new products SO I can expand my basket", CYAN, False, 9),
    ("  without extra effort or delivery charges.", CYAN, False, 9),
    ("", WHITE),
    ("How (from NextLeap + Swiggy Instamart insights):", WHITE, True, 10),
    ("  1. 'Try Something New' section: 1-2 tail-category", GRAY, False, 9),
    ("     items added to cart suggestions at checkout", GRAY, False, 9),
    ("  2. Smart Basket: AI restocking + new item sampling", GRAY, False, 9),
    ("  3. Category-specific fulfillment (Swiggy model) for", GRAY, False, 9),
    ("     wider tail catalog without slowing head delivery", GRAY, False, 9),
    ("  4. Gamified exploration: 'Explore 3 new categories", GRAY, False, 9),
    ("     this month -> unlock free delivery credits'", GRAY, False, 9),
    ("", WHITE),
    ("Expected Impact:", YELLOW, True, 10),
    ("  Reverse Medium Exploration decline (-4.8% -> +5%)", YELLOW, False, 9),
    ("  Increase AOV by 15-20% via cross-category discovery", YELLOW, False, 9),
    ("  Convert Low Exploration stable users into Medium", YELLOW, False, 9),
])

# ── BOTTOM: Prioritization Matrix ────────────────────────────────────
_rect(s2, 0.3, 5.4, 12.7, 0.02, DGRAY)

_card(s2, 0.3, 5.6, 8.5, 1.7, CYAN)
_ml(s2, 0.5, 5.65, 8.2, 1.6, [
    ("PRIORITIZATION ROADMAP", CYAN, True, 12),
    ("", WHITE),
    ("Q1 (Immediate):  Solution 2 -- Instant Resolution Engine (low effort, high impact on rating)", WHITE, False, 10),
    ("  Quick win: auto-refund + extended returns. Directly improves 3.7 -> 4.0 rating.", GRAY, False, 9),
    ("Q2 (Short-term):  Solution 1 -- Predictive ETA (medium effort, critical severity)", WHITE, False, 10),
    ("  Addresses #1 barrier. Requires ML model + dark store instrumentation.", GRAY, False, 9),
    ("Q3 (Medium-term): Solution 3 -- AI Discovery Engine (higher effort, long-term AOV growth)", WHITE, False, 10),
    ("  Drives category exploration + retention. Needs recommendation engine + UX redesign.", GRAY, False, 9),
])

_card(s2, 9.0, 5.6, 4.0, 1.7, YELLOW)
_ml(s2, 9.2, 5.65, 3.6, 1.6, [
    ("SUCCESS METRICS", YELLOW, True, 12),
    ("", WHITE),
    ("Avg Rating: 3.7 -> 4.2+ (2 quarters)", WHITE, False, 10),
    ("Delivery Complaints: -40%", GREEN, False, 10),
    ("Refund Resolution Time: <2 min", GREEN, False, 10),
    ("Medium Exploration: -4.8% -> +5%", GREEN, False, 10),
    ("AOV Increase: +15-20%", GREEN, False, 10),
    ("Support Gap: 58% -> 72% (parity)", GREEN, False, 10),
])

# ── save ──────────────────────────────────────────────────────────────
out = r"c:\Top fellow - blinkit\Blinkit_Problem_Report.pptx"
prs.save(out)
print(f"Saved to: {out}")
print(f"Total slides: {len(prs.slides)}")
